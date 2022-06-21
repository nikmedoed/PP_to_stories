import os
import sys
import time

import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SETTINGS_EXPORT = {
    'video': {'VertResolution': 1920, 'DefaultSlideDuration': 1},
    'gif': {'VertResolution': 1920, 'DefaultSlideDuration': 15},
}


def media_type_det(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
        return "video"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            if hasattr(shape, 'image') and shape.image.content_type == 'image/gif':
                return 'gif'
        except:
            pass
    return None


def slide_type_det(slide):
    slide_types = set()
    for shape in slide.shapes:
        slide_types.add(media_type_det(shape))
    for slide_type in SETTINGS_EXPORT.keys():
        if slide_type in slide_types:
            break
    else:
        slide_type = None
    return slide_type


def pptx_to_stories(pptx_path):
    result_dir = f'{os.path.splitext(pptx_path)[0]}-slides'
    pptx_name = os.path.basename(result_dir)

    prs = Presentation(pptx_path)
    save_as_video = [slide_type_det(slide) for slide in prs.slides]

    try:
        os.mkdir(result_dir)
    except:
        pass

    Application = win32com.client.Dispatch("PowerPoint.Application")

    read_only = True
    has_title = False
    window = False
    slides = 0
    slides_total = len(save_as_video)

    if not all(save_as_video):
        presentation = Application.Presentations.Open(pptx_path, read_only, has_title, window)
        for n, as_video in enumerate(save_as_video):
            if not as_video:
                presentation.Slides[n].Export(f"{result_dir}\\{n:02}-{pptx_name}.png", "PNG")
                slides += 1
        presentation.Close()

    print(f'completed: {slides} / {slides_total}')

    for n, as_video in enumerate(save_as_video):
        presentation = Application.Presentations.Open(pptx_path, read_only, has_title, window)
        if as_video:
            for i in range(len(save_as_video) - 1, -1, -1):
                if i != n:
                    presentation.Slides[i].Delete()
            presentation.CreateVideo(f"{result_dir}\\{n:02}-{pptx_name}.mp4", **SETTINGS_EXPORT[as_video])
            time.sleep(2)
            while presentation.CreateVideoStatus == 1:
                time.sleep(1)
            slides += 1
            print(f'completed: {slides} / {slides_total}')
        presentation.Close()


if __name__ == "__main__":
    if len(sys.argv) > 1:
        print(sys.argv)
        pptx_path = sys.argv[1]
    else:
        print("example")
        pptx_path = f'{os.getcwd()}\\example\\example2.pptx'
    pptx_to_stories(pptx_path)
    print("ok")
